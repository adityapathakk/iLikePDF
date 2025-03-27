"""
this python script is an end-to-end pdf-to-pptx converter.
instructions to use:
- using requirements.txt, install all libraries
- in the directory where this script is present, create a directory - `pdf-pptx`
- in `pdf-pptx`, create a directory - `to-convert`
- transfer all the pdfs that need to be converted to this directory (i.e. `pdf-pptx/to-convert`)
- run this script! your outputs will be organised in `pdf-pptx/working/{pdf_name}`
"""


### importing libraries
import os, re, gc, time
from pptx import Presentation as pptx_pres
from spire.pdf.common import *
from spire.pdf import *
from spire.pdf import FileFormat as spirePDF_FF
from spire.presentation import *
from spire.presentation import FileFormat as FF
from spire.presentation.common import *

### step 1 - splitting PDFs
def split_pdf(pdf_path, pdf_name):
    # creating a directory to store the split PDFs
    os.makedirs(f"pdf-pptx/working/{pdf_name}/split_pdfs", exist_ok = True)
    # creating a PdfDocument object
    target_doc = PdfDocument()
    # loading target PDF file
    target_doc.LoadFromFile(pdf_path)

    # creating PdfDocument objects that will consist of the split target PDF
    num_required_splits = target_doc.Pages.Count // 10 + 1
    docs = [PdfDocument() for _ in range(num_required_splits)]
    
    doc_index = 0
    remainingPages = target_doc.Pages.Count
    current_page = 0
    while remainingPages >= 1:
        if remainingPages >= 10:
            docs[doc_index].InsertPageRange(target_doc, current_page, current_page + 9)
            current_page += 10
            remainingPages -= 10
            doc_index += 1
            continue
        elif remainingPages >= 1:
            docs[doc_index].InsertPageRange(target_doc, current_page, current_page + remainingPages - 1)
            doc_index += 1
            remainingPages = 0
            current_page += remainingPages
            break
    
    for i, doc in enumerate(docs):
        doc.SaveToFile(f"pdf-pptx/working/{pdf_name}/split_pdfs/Split-{i + 1}.pdf")
        doc.Close() # closing the PdfDocument object

    target_doc.Close()
    return f"pdf-pptx/working/{pdf_name}/split_pdfs" # return the path to the directory containing the split PDFs

### step 2 - converting split PDFs to PPTXs
def convert_pdf(split_pdf_dir, pdf_name):
    os.makedirs(f"pdf-pptx/working/{pdf_name}/split_pptxs", exist_ok = True)
    print("Directory containing split PDFs: ", split_pdf_dir)
    for split_pdf_path in os.listdir(split_pdf_dir):
        print(f"Converting {split_pdf_path} to PPTX...")
        pdf = PdfDocument()
        try: 
            pdf.LoadFromFile(os.path.join(split_pdf_dir, split_pdf_path))
            print(f"Loaded {os.path.join(split_pdf_dir, split_pdf_path)}")
            # converting split PDF file to PPTX file
            pdf.SaveToFile(f"pdf-pptx/working/{pdf_name}/split_pptxs/{split_pdf_path.split(".")[0]}.pptx", spirePDF_FF.PPTX)
        finally:
            pdf.Close()
            print(f"Closed {split_pdf_path}")
    
    return f"pdf-pptx/working/{pdf_name}/split_pptxs" # return the path to the directory containing the split PPTXs

### step 3 - merging split PPTXs into one PPTX
def merge_pptxs(split_pptx_dir, pdf_name):
    os.makedirs(f"pdf-pptx/working/{pdf_name}/temp-output", exist_ok = True)
    target_pres = Presentation()
    target_pres.LoadFromFile(f"pdf-pptx/working/{pdf_name}/split_pptxs/Split-1.pptx")
    
    for i in range(2, len(os.listdir(split_pptx_dir)) + 1):
        pres = Presentation()
        pres.LoadFromFile(f"pdf-pptx/working/{pdf_name}/split_pptxs/Split-{i}.pptx")
        for slide in pres.Slides:
            target_pres.Slides.AppendBySlide(slide)
        pres.Dispose()
    
    target_pres.SaveToFile(f"pdf-pptx/working/{pdf_name}/temp-output/{pdf_name}-spire.pptx", FF.Pptx2016)
    target_pres.Dispose()
    return f"pdf-pptx/working/{pdf_name}/temp-output/{pdf_name}-spire.pptx" # return the path to the merged PPTX with spire watermarks

### step 4 - removing spire watermarks
def remove_watermarks(spire_pptx_path, pdf_name):
    # defining regex patterns for the two evaluation warnings
    eval1 = re.compile("Evaluation Warning : The document was created with Spire.PDF for Python.")
    eval2 = re.compile("Evaluation Warning : The document was created with Spire.Presentation for Python")

    def remove_shapes_with_regex(slide, regex):
        """Remove shapes in the slide where the regex pattern is found in the text."""
        shapes_to_remove = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                # Check if any paragraph in the shape contains the pattern
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if regex.search(run.text):
                            shapes_to_remove.append(shape)
                            break
                    else:
                        continue
                    break
        
        # Remove the shapes that contain the regex pattern
        for shape in shapes_to_remove:
            slide.shapes._spTree.remove(shape._element)

    def presentation_remove_shapes(presentation, regex):
        """Remove shapes containing the regex pattern in all slides of the presentation."""
        for slide in presentation.slides:
            remove_shapes_with_regex(slide, regex)
        return presentation

    pres = pptx_pres(spire_pptx_path)

    presentation_remove_shapes(pres, eval1)
    presentation_remove_shapes(pres, eval2)

    pres.save(f"pdf-pptx/working/{pdf_name}/{pdf_name}.pptx")

### main
workingDir = "pdf-pptx/working"
os.makedirs(workingDir, exist_ok = True)

# path to the directory containing PDF files to be converted
path = "pdf-pptx/to-convert"
for i, pdf in enumerate(os.listdir(path)):
    if i == 0:
        continue
    pdf_name = pdf.split(".")[0]
    pdf_path = f"{path}/{pdf}"
    print(f"Splitting ({pdf_name}) at {pdf_path}...")
    split_pdf_dir = split_pdf(pdf_path, pdf_name)
    print(f"({pdf_name}) split into multiple PDFs. Converting to PPTX...")
    split_pptx_dir = convert_pdf(split_pdf_dir, pdf_name)
    print(f"({pdf_name})'s splits converted into PPTX. Merging into one PPTX...")
    spire_pptx_path = merge_pptxs(split_pptx_dir, pdf_name)
    print(f"({spire_pptx_path.split(".")[0]}) created. Removing watermarks...")
    remove_watermarks(spire_pptx_path, pdf_name)
    print(f"({pdf_name}) converted to PPTX!\n")

    # collected = gc.collect() # garbage collection to free up memory
    # print(f"Garbage collected: {collected}\n")
    # time.sleep(5) # sleep for 5 seconds to avoid memory issues